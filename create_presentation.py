"""Create Mood Weather presentation (5 slides)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ──
BG_DARK     = RGBColor(0x1B, 0x2A, 0x4A)  # deep navy
BG_ACCENT   = RGBColor(0x2C, 0x3E, 0x6B)  # lighter navy
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCC, 0xD6, 0xE0)
ACCENT_BLUE = RGBColor(0x42, 0x99, 0xE1)
ACCENT_GOLD = RGBColor(0xF6, 0xE0, 0x5E)
ACCENT_GREEN= RGBColor(0x48, 0xBB, 0x78)
ACCENT_PINK = RGBColor(0xF6, 0x87, 0xB3)
DARK_TEXT     = RGBColor(0x2D, 0x37, 0x48)
CARD_BG     = RGBColor(0x23, 0x35, 0x5A)

# ── helpers ──
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=LIGHT_GRAY, spacing=Pt(8), bullet_char="▸"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char}  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
    return txBox

def add_card(slide, left, top, width, height, title, items, accent=ACCENT_BLUE):
    """Add a styled card with title and bullet items."""
    card = add_shape(slide, left, top, width, height, CARD_BG, accent)
    # title bar
    add_textbox(slide, left + Inches(0.3), top + Inches(0.15), width - Inches(0.6), Inches(0.5),
                title, font_size=18, color=accent, bold=True)
    # thin line under title
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   left + Inches(0.3), top + Inches(0.6),
                                   width - Inches(0.6), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = accent
    line.line.fill.background()
    # bullets
    add_bullet_list(slide, left + Inches(0.3), top + Inches(0.75),
                    width - Inches(0.6), height - Inches(0.9),
                    items, font_size=14, color=LIGHT_GRAY)

# ================================================================
# SLIDE 1 — Title
# ================================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide1, BG_DARK)

# Decorative top bar
bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_GOLD; bar.line.fill.background()

# Emoji + title
add_textbox(slide1, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
            "☀️  Mood Weather", font_size=52, color=WHITE, bold=True, font_name="Calibri")

add_textbox(slide1, Inches(1.5), Inches(2.8), Inches(10), Inches(0.6),
            "A Visual Mood Tracker Using Weather Metaphors",
            font_size=26, color=ACCENT_GOLD, bold=False)

# Divider
div = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(3.7), Inches(4), Pt(3))
div.fill.solid(); div.fill.fore_color.rgb = ACCENT_BLUE; div.line.fill.background()

# Info lines
info_lines = [
    "Your Name  ·  Your University Email  ·  Your Group",
    "",
    "Software Engineering Toolkit  —  Final Project Presentation",
]
add_bullet_list(slide1, Inches(1.5), Inches(4.0), Inches(10), Inches(2),
                info_lines, font_size=18, color=LIGHT_GRAY, spacing=Pt(6), bullet_char="")

# ================================================================
# SLIDE 2 — Context
# ================================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2, BG_DARK)
bar2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar2.fill.solid(); bar2.fill.fore_color.rgb = ACCENT_GOLD; bar2.line.fill.background()

add_textbox(slide2, Inches(0.8), Inches(0.3), Inches(12), Inches(0.7),
            "Context", font_size=38, color=WHITE, bold=True)
add_textbox(slide2, Inches(0.8), Inches(1.0), Inches(12), Inches(0.4),
            "Who, Why, and What", font_size=20, color=ACCENT_BLUE)

# Three cards
add_card(slide2, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8),
         "👤  End User",
         [
             "People who want to track their daily mood",
             "Prefer simple, visual interfaces",
             "All ages — no technical skills needed",
         ], accent=ACCENT_BLUE)

add_card(slide2, Inches(4.9), Inches(1.8), Inches(3.6), Inches(4.8),
         "💡  Problem",
         [
             "Traditional trackers are tedious",
             "Long questionnaires → users quit",
             "Rating scales feel unnatural",
             "Lack of visual engagement",
         ], accent=ACCENT_PINK)

add_card(slide2, Inches(9.0), Inches(1.8), Inches(3.6), Inches(4.8),
         "🎯  Product Idea",
         [
             "One tap per day",
             "13 weather icons = 13 moods",
             "Personal weather calendar",
             "Easy to understand at a glance",
         ], accent=ACCENT_GREEN)

# ================================================================
# SLIDE 3 — Implementation
# ================================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3, BG_DARK)
bar3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar3.fill.solid(); bar3.fill.fore_color.rgb = ACCENT_GOLD; bar3.line.fill.background()

add_textbox(slide3, Inches(0.8), Inches(0.3), Inches(12), Inches(0.7),
            "Implementation", font_size=38, color=WHITE, bold=True)

# Tech stack row
add_textbox(slide3, Inches(0.8), Inches(1.1), Inches(12), Inches(0.4),
            "Tech Stack", font_size=22, color=ACCENT_GOLD, bold=True)

tech_items = [
    "Backend:  Python 3.12 · FastAPI · SQLModel · asyncpg · PostgreSQL 18",
    "Frontend: React 18 · TypeScript · Vite · Chart.js",
    "Proxy:    Caddy 2 — reverse proxy",
    "Infra:    Docker Compose — full container orchestration",
    "Deploy:   Ubuntu 24.04 VM — automated via SSH + Docker Compose",
    "Testing:  pytest · pytest-asyncio · httpx",
]
add_bullet_list(slide3, Inches(0.8), Inches(1.55), Inches(11.5), Inches(2.2),
                tech_items, font_size=15, color=LIGHT_GRAY, spacing=Pt(5), bullet_char="•")

# V1 vs V2
add_textbox(slide3, Inches(0.8), Inches(3.9), Inches(12), Inches(0.4),
            "Version 1  →  Version 2", font_size=22, color=ACCENT_GOLD, bold=True)

# V1 card
add_card(slide3, Inches(0.8), Inches(4.5), Inches(5.5), Inches(2.7),
         "Version 1 — Core",
         [
             "CRUD API (create, list, delete moods)",
             "4 weather types (sunny, cloudy, rainy, stormy)",
             "Basic React form for mood entry",
             "Docker Compose local development",
         ], accent=ACCENT_BLUE)

# V2 card
add_card(slide3, Inches(6.9), Inches(4.5), Inches(5.6), Inches(2.7),
         "Version 2 — Enhanced",
         [
             "13 weather types for richer mood expression",
             "PUT endpoint — edit existing mood entries",
             "Inline edit & delete on frontend",
             "Monthly calendar view with weather icons",
             "Analytics dashboard (Chart.js)",
             "Deployed to production VM",
         ], accent=ACCENT_GREEN)

# ================================================================
# SLIDE 4 — Demo
# ================================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide4, BG_DARK)
bar4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar4.fill.solid(); bar4.fill.fore_color.rgb = ACCENT_GOLD; bar4.line.fill.background()

add_textbox(slide4, Inches(0.8), Inches(0.3), Inches(12), Inches(0.7),
            "Demo — Version 2", font_size=38, color=WHITE, bold=True)
add_textbox(slide4, Inches(0.8), Inches(1.0), Inches(12), Inches(0.4),
            "Pre-recorded Video Demonstration with Voice-Over  (≤ 2 min)",
            font_size=20, color=ACCENT_BLUE)

# Placeholder box for video
video_box = add_shape(slide4, Inches(1.5), Inches(1.8), Inches(10.3), Inches(5.2),
                      CARD_BG, ACCENT_GOLD)
video_box.text_frame.word_wrap = True
tf = video_box.text_frame
p = tf.paragraphs[0]
p.text = "🎬  Insert your demo video here"
p.font.size = Pt(28)
p.font.color.rgb = ACCENT_GOLD
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
tf.paragraphs[0].space_before = Pt(80)

p2 = tf.add_paragraph()
p2.text = "(Drag & drop your .mp4 onto this slide in PowerPoint)"
p2.font.size = Pt(16)
p2.font.color.rgb = LIGHT_GRAY
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(12)

p3 = tf.add_paragraph()
p3.text = "What to show: login → add mood → edit entry → delete entry → calendar view → dashboard"
p3.font.size = Pt(14)
p3.font.color.rgb = ACCENT_BLUE
p3.alignment = PP_ALIGN.CENTER
p3.space_before = Pt(30)

# ================================================================
# SLIDE 5 — Links
# ================================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide5, BG_DARK)
bar5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar5.fill.solid(); bar5.fill.fore_color.rgb = ACCENT_GOLD; bar5.line.fill.background()

add_textbox(slide5, Inches(0.8), Inches(0.3), Inches(12), Inches(0.7),
            "Links", font_size=38, color=WHITE, bold=True)
add_textbox(slide5, Inches(0.8), Inches(1.0), Inches(12), Inches(0.4),
            "Access the project online", font_size=20, color=ACCENT_BLUE)

# GitHub card
add_card(slide5, Inches(1.5), Inches(1.8), Inches(4.8), Inches(2.5),
         "📦  GitHub Repository",
         [
             "https://github.com/your-username/se-toolkit-mood-weather",
             "Full source code (backend + frontend)",
             "Docker configuration & CI/CD",
         ], accent=ACCENT_GOLD)

# Deployed product card
add_card(slide5, Inches(7.0), Inches(1.8), Inches(4.8), Inches(2.5),
         "🚀  Deployed Product",
         [
             "http://<your-vm-ip>:42002  —  Frontend",
             "http://<your-vm-ip>:42001/docs  —  API Docs",
             "Live & accessible 24/7",
         ], accent=ACCENT_GREEN)

# QR code placeholders
add_textbox(slide5, Inches(0.8), Inches(4.6), Inches(12), Inches(0.4),
            "QR Codes", font_size=22, color=ACCENT_GOLD, bold=True)

qr1 = add_shape(slide5, Inches(1.5), Inches(5.2), Inches(2.5), Inches(2.0), WHITE)
qr1.text_frame.word_wrap = True
tf_qr = qr1.text_frame
tf_qr.word_wrap = True
p = tf_qr.paragraphs[0]
p.text = "📱 GitHub QR"
p.font.size = Pt(14)
p.font.color.rgb = DARK_TEXT
p.alignment = PP_ALIGN.CENTER
tf_qr.paragraphs[0].space_before = Pt(30)

qr2 = add_shape(slide5, Inches(5.5), Inches(5.2), Inches(2.5), Inches(2.0), WHITE)
qr2.text_frame.word_wrap = True
tf_qr2 = qr2.text_frame
tf_qr2.word_wrap = True
p = tf_qr2.paragraphs[0]
p.text = "🌐 App QR"
p.font.size = Pt(14)
p.font.color.rgb = DARK_TEXT
p.alignment = PP_ALIGN.CENTER
tf_qr2.paragraphs[0].space_before = Pt(30)

# Thanks
add_textbox(slide5, Inches(8.5), Inches(5.2), Inches(4), Inches(2),
            "Thank you! 🙏",
            font_size=32, color=ACCENT_GOLD, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide5, Inches(8.5), Inches(6.0), Inches(4), Inches(0.5),
            "Questions?",
            font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ── Save ──
output_path = r"c:\Users\lis_b\Desktop\software-engineering-toolkit\se-toolkit-mood-weather\Mood_Weather_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
